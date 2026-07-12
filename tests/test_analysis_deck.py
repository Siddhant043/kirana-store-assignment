"""Analysis deck generation tests against Postgres."""

from datetime import datetime, timedelta
from decimal import Decimal
from io import BytesIO

import pytest
from pptx import Presentation
from sqlalchemy.ext.asyncio import AsyncSession

from src.db.models import Bill, BillLine, DraftBill, Product
from src.domain.analysis_deck import (
    REQUIRED_SLIDE_TITLES,
    SLIDE_GST_BY_SLAB,
    SLIDE_INSIGHTS,
    SLIDE_PAYMENT_MIX,
    SLIDE_SALES_TREND,
    SLIDE_STOCK_HEALTH,
    SLIDE_TOP_ITEMS,
    compose_insights,
    gather_analysis_deck_data,
    generate_analysis_deck,
)
from src.domain.analytics import AnalyticsService
from src.domain.doc_gen import run_cpu_bound
from src.domain.inventory import InventoryService
from src.domain.invoice import format_paise_as_rupees
from src.domain.shop_time import today_ist, utc_bounds_for_ist_date

CHAT_ID = 77
_invoice_counter = 0


def _slide_texts(slide: object) -> str:
    parts: list[str] = []
    for shape in getattr(slide, "shapes"):
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _find_slide(presentation: Presentation, title: str) -> object:
    for slide in presentation.slides:
        if title in _slide_texts(slide):
            return slide
    msg = f"slide titled {title!r} not found"
    raise AssertionError(msg)


def _first_chart(slide: object) -> object:
    for shape in getattr(slide, "shapes"):
        if getattr(shape, "has_chart", False):
            return shape.chart
    msg = "no chart on slide"
    raise AssertionError(msg)


def _series_values(chart: object, series_index: int = 0) -> list[float]:
    plot = chart.plots[0]
    return [float(value) for value in plot.series[series_index].values]


async def _create_product(
    session: AsyncSession,
    *,
    name: str,
    mrp_paise: int,
    gst_slab: int,
    hsn_code: str,
    quantity: Decimal,
    reorder_level: Decimal,
) -> Product:
    service = InventoryService(session)
    result = await service.add_product(
        name=name,
        brand=None,
        mrp_paise=mrp_paise,
        cost_price_paise=mrp_paise - 100,
        gst_slab=gst_slab,
        hsn_code=hsn_code,
        unit_type="packaged",
        reorder_level=reorder_level,
    )
    assert result.status == "ok"
    product = await session.get(Product, result.product_id)
    assert product is not None
    product.quantity = quantity
    await session.flush()
    return product


async def _seed_bill(
    session: AsyncSession,
    *,
    finalized_at: datetime,
    payment_mode: str,
    lines: list[tuple[Product, Decimal, int, int, int, int]],
    total_paise: int,
    cgst_paise: int,
    sgst_paise: int,
) -> Bill:
    global _invoice_counter
    _invoice_counter += 1

    draft = DraftBill(chat_id=CHAT_ID, status="finalized")
    session.add(draft)
    await session.flush()

    bill = Bill(
        draft_bill_id=draft.draft_bill_id,
        chat_id=CHAT_ID,
        invoice_number=f"INV-DECK-{_invoice_counter:04d}",
        payment_mode=payment_mode,
        subtotal_paise=total_paise - cgst_paise - sgst_paise,
        cgst_paise=cgst_paise,
        sgst_paise=sgst_paise,
        round_off_paise=0,
        total_paise=total_paise,
        finalized_at=finalized_at,
    )
    session.add(bill)
    await session.flush()
    draft.bill_id = bill.bill_id

    for product, quantity, line_total, taxable, line_cgst, line_sgst in lines:
        session.add(
            BillLine(
                bill_id=bill.bill_id,
                product_id=product.product_id,
                quantity=quantity,
                mrp_paise=product.mrp_paise,
                cost_price_paise=product.cost_price_paise,
                gst_slab=product.gst_slab,
                hsn_code=product.hsn_code,
                line_total_paise=line_total,
                taxable_paise=taxable,
                cgst_paise=line_cgst,
                sgst_paise=line_sgst,
            )
        )
    await session.flush()
    return bill


@pytest.mark.asyncio
async def test_analysis_deck_has_required_slides_and_matching_charts(
    inventory_session: AsyncSession,
) -> None:
    today = today_ist()
    start_utc, _ = utc_bounds_for_ist_date(today - timedelta(days=1))
    today_start, _ = utc_bounds_for_ist_date(today)

    atta = await _create_product(
        inventory_session,
        name="Deck Atta",
        mrp_paise=10000,
        gst_slab=5,
        hsn_code="110100",
        quantity=Decimal("2"),
        reorder_level=Decimal("10"),
    )
    maggi = await _create_product(
        inventory_session,
        name="Deck Maggi",
        mrp_paise=2000,
        gst_slab=12,
        hsn_code="190230",
        quantity=Decimal("50"),
        reorder_level=Decimal("5"),
    )

    await _seed_bill(
        inventory_session,
        finalized_at=start_utc + timedelta(hours=2),
        payment_mode="cash",
        lines=[(atta, Decimal("1"), 10000, 9524, 238, 238)],
        total_paise=10000,
        cgst_paise=238,
        sgst_paise=238,
    )
    await _seed_bill(
        inventory_session,
        finalized_at=today_start + timedelta(hours=3),
        payment_mode="upi",
        lines=[
            (maggi, Decimal("3"), 6000, 5357, 321, 322),
            (atta, Decimal("1"), 10000, 9524, 238, 238),
        ],
        total_paise=16000,
        cgst_paise=559,
        sgst_paise=560,
    )

    deck_data = await gather_analysis_deck_data(inventory_session, day_count=7)
    insights = compose_insights(deck_data)
    pptx_bytes = generate_analysis_deck(deck_data, insights)
    presentation = Presentation(BytesIO(pptx_bytes))

    titles_found = {_slide_texts(slide) for slide in presentation.slides}
    for required in REQUIRED_SLIDE_TITLES:
        assert any(required in text for text in titles_found), required

    analytics = AnalyticsService(inventory_session)
    trend = await analytics.sales_trend(7)
    payment = deck_data.payment_mode_split
    gst_rows = await analytics.gst_collected_by_slab()

    trend_slide = _find_slide(presentation, SLIDE_SALES_TREND)
    assert _series_values(_first_chart(trend_slide)) == [
        float(point.total_sales_paise) for point in trend
    ]

    payment_slide = _find_slide(presentation, SLIDE_PAYMENT_MIX)
    assert _series_values(_first_chart(payment_slide)) == [
        float(payment.cash_paise),
        float(payment.upi_paise),
        float(payment.card_paise),
        float(payment.khata_paise),
    ]

    gst_slide = _find_slide(presentation, SLIDE_GST_BY_SLAB)
    assert _series_values(_first_chart(gst_slide)) == [
        float(row.tax_collected_paise) for row in gst_rows
    ]

    top_slide = _find_slide(presentation, SLIDE_TOP_ITEMS)
    charts = [shape.chart for shape in top_slide.shapes if shape.has_chart]
    assert len(charts) == 2
    assert _series_values(charts[0]) == [
        float(item.revenue_paise) for item in deck_data.top_items_by_revenue[:5]
    ]

    stock_slide = _find_slide(presentation, SLIDE_STOCK_HEALTH)
    stock_chart = _first_chart(stock_slide)
    categories = [str(category.label) for category in stock_chart.plots[0].categories]
    assert "Deck Atta" in categories
    atta_index = categories.index("Deck Atta")
    assert _series_values(stock_chart, 0)[atta_index] == 2.0
    assert _series_values(stock_chart, 1)[atta_index] == 10.0


@pytest.mark.asyncio
async def test_compose_insights_only_uses_analytics_figures(
    inventory_session: AsyncSession,
) -> None:
    today = today_ist()
    today_start, _ = utc_bounds_for_ist_date(today)
    product = await _create_product(
        inventory_session,
        name="Insight Oil",
        mrp_paise=5000,
        gst_slab=5,
        hsn_code="1512",
        quantity=Decimal("100"),
        reorder_level=Decimal("1"),
    )
    await _seed_bill(
        inventory_session,
        finalized_at=today_start + timedelta(hours=1),
        payment_mode="khata",
        lines=[(product, Decimal("2"), 10000, 9524, 238, 238)],
        total_paise=10000,
        cgst_paise=238,
        sgst_paise=238,
    )

    deck_data = await gather_analysis_deck_data(inventory_session, day_count=7)
    insights = compose_insights(deck_data)

    assert format_paise_as_rupees(deck_data.total_sales_paise) in insights
    assert str(deck_data.bill_count) in insights
    assert format_paise_as_rupees(deck_data.tax_collected_paise) in insights
    assert format_paise_as_rupees(deck_data.payment_mode_split.khata_paise) in insights
    assert "Insight Oil" in insights
    assert "estimated" not in insights.lower()
    assert "probably" not in insights.lower()

    pptx_bytes = generate_analysis_deck(deck_data, insights)
    presentation = Presentation(BytesIO(pptx_bytes))
    insights_slide = _find_slide(presentation, SLIDE_INSIGHTS)
    assert format_paise_as_rupees(deck_data.total_sales_paise) in _slide_texts(
        insights_slide
    )


@pytest.mark.asyncio
async def test_generate_analysis_deck_via_doc_gen_lock(
    inventory_session: AsyncSession,
) -> None:
    deck_data = await gather_analysis_deck_data(inventory_session, day_count=7)
    insights = compose_insights(deck_data)
    pptx_bytes = await run_cpu_bound(generate_analysis_deck, deck_data, insights)
    presentation = Presentation(BytesIO(pptx_bytes))
    assert len(presentation.slides) == len(REQUIRED_SLIDE_TITLES)
