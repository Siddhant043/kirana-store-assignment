"""On-request analysis deck (python-pptx) from analytics figures."""

from dataclasses import dataclass
from decimal import Decimal
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from sqlalchemy.ext.asyncio import AsyncSession

from src.domain.analytics import (
    AnalyticsService,
    DailySalesPoint,
    GstSlabTotal,
    PaymentModeSplit,
    TopItem,
)
from src.domain.inventory import InventoryService, LowStockProduct
from src.domain.invoice import format_paise_as_rupees
from src.domain.shop_time import rolling_last_n_ist_days

SLIDE_TITLE = "Store Analysis"
SLIDE_SALES_TREND = "Sales Trend"
SLIDE_TOP_ITEMS = "Top Items"
SLIDE_PAYMENT_MIX = "Payment Mix"
SLIDE_GST_BY_SLAB = "GST Collected by GST Slab"
SLIDE_STOCK_HEALTH = "Stock Health"
SLIDE_INSIGHTS = "Insights"

REQUIRED_SLIDE_TITLES = (
    SLIDE_TITLE,
    SLIDE_SALES_TREND,
    SLIDE_TOP_ITEMS,
    SLIDE_PAYMENT_MIX,
    SLIDE_GST_BY_SLAB,
    SLIDE_STOCK_HEALTH,
    SLIDE_INSIGHTS,
)


@dataclass(frozen=True)
class AnalysisDeckData:
    period_start: str
    period_end: str
    bill_count: int
    total_sales_paise: int
    tax_collected_paise: int
    sales_trend: list[DailySalesPoint]
    top_items_by_revenue: list[TopItem]
    top_items_by_quantity: list[TopItem]
    payment_mode_split: PaymentModeSplit
    gst_by_slab: list[GstSlabTotal]
    low_stock: list[LowStockProduct]


async def gather_analysis_deck_data(
    session: AsyncSession,
    *,
    day_count: int = 7,
) -> AnalysisDeckData:
    analytics = AnalyticsService(session)
    inventory = InventoryService(session)
    period_start, period_end = rolling_last_n_ist_days(day_count)
    report = await analytics.report_for_ist_range(period_start, period_end)
    sales_trend = await analytics.sales_trend(day_count)
    gst_by_slab = await analytics.gst_collected_by_slab(period_start, period_end)
    top_by_qty = await analytics.top_items_by_quantity(period_start, period_end)
    low_stock = await inventory.list_low_stock()
    return AnalysisDeckData(
        period_start=period_start.isoformat(),
        period_end=period_end.isoformat(),
        bill_count=report.bill_count,
        total_sales_paise=report.total_sales_paise,
        tax_collected_paise=report.tax_collected_paise,
        sales_trend=sales_trend,
        top_items_by_revenue=report.top_items,
        top_items_by_quantity=top_by_qty,
        payment_mode_split=report.payment_mode_split,
        gst_by_slab=gst_by_slab,
        low_stock=low_stock.products,
    )


def compose_insights(data: AnalysisDeckData) -> str:
    payment = data.payment_mode_split
    lines = [
        f"Period {data.period_start} to {data.period_end} (IST).",
        (
            f"Total sales ₹{format_paise_as_rupees(data.total_sales_paise)} "
            f"across {data.bill_count} Bills."
        ),
        (
            f"Tax collected ₹{format_paise_as_rupees(data.tax_collected_paise)} "
            f"(CGST+SGST from finalized Bills)."
        ),
        (
            "Payment Mode split — "
            f"cash ₹{format_paise_as_rupees(payment.cash_paise)}, "
            f"UPI ₹{format_paise_as_rupees(payment.upi_paise)}, "
            f"card ₹{format_paise_as_rupees(payment.card_paise)}, "
            f"Khata ₹{format_paise_as_rupees(payment.khata_paise)}."
        ),
    ]
    if data.top_items_by_revenue:
        top = data.top_items_by_revenue[0]
        lines.append(
            f"Top item by revenue: {top.name} "
            f"(₹{format_paise_as_rupees(top.revenue_paise)})."
        )
    if data.gst_by_slab:
        slab_bits = ", ".join(
            f"{row.gst_slab}% ₹{format_paise_as_rupees(row.tax_collected_paise)}"
            for row in data.gst_by_slab
        )
        lines.append(f"GST by GST Slab: {slab_bits}.")
    lines.append(f"Products at/below Reorder Level: {len(data.low_stock)}.")
    return "\n".join(lines)


def generate_analysis_deck(data: AnalysisDeckData, insights_text: str) -> bytes:
    presentation = Presentation()
    _add_title_slide(presentation, data)
    _add_sales_trend_slide(presentation, data)
    _add_top_items_slide(presentation, data)
    _add_payment_mix_slide(presentation, data)
    _add_gst_by_slab_slide(presentation, data)
    _add_stock_health_slide(presentation, data)
    _add_insights_slide(presentation, insights_text)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _add_title_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_TITLE)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8), Inches(2))
    text_frame = body.text_frame
    text_frame.text = (
        f"IST period {data.period_start} to {data.period_end}\n"
        f"Bills: {data.bill_count}\n"
        f"Total sales: ₹{format_paise_as_rupees(data.total_sales_paise)}"
    )


def _add_sales_trend_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_SALES_TREND)
    chart_data = _new_category_chart_data()
    categories = [point.business_date for point in data.sales_trend] or ["(none)"]
    chart_data.categories = categories
    values = (
        [point.total_sales_paise for point in data.sales_trend]
        if data.sales_trend
        else [0]
    )
    _add_chart_series(chart_data, "Sales (paise)", values)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5),
        chart_data,
    )


def _add_top_items_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_TOP_ITEMS)

    revenue_items = data.top_items_by_revenue[:5] or [
        TopItem(product_id=0, name="(none)", quantity="0", revenue_paise=0)
    ]
    chart_data = _new_category_chart_data()
    chart_data.categories = [item.name for item in revenue_items]
    _add_chart_series(
        chart_data,
        "Revenue (paise)",
        [item.revenue_paise for item in revenue_items],
    )
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4),
        Inches(1.4),
        Inches(4.5),
        Inches(5),
        chart_data,
    )

    qty_items = data.top_items_by_quantity[:5] or [
        TopItem(product_id=0, name="(none)", quantity="0", revenue_paise=0)
    ]
    qty_chart = _new_category_chart_data()
    qty_chart.categories = [item.name for item in qty_items]
    _add_chart_series(
        qty_chart,
        "Quantity",
        [_quantity_as_float(item.quantity) for item in qty_items],
    )
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(5.1),
        Inches(1.4),
        Inches(4.5),
        Inches(5),
        qty_chart,
    )


def _add_payment_mix_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_PAYMENT_MIX)
    split = data.payment_mode_split
    chart_data = _new_category_chart_data()
    chart_data.categories = ["cash", "upi", "card", "khata"]
    _add_chart_series(
        chart_data,
        "Payment Mode (paise)",
        [
            split.cash_paise,
            split.upi_paise,
            split.card_paise,
            split.khata_paise,
        ],
    )
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1.5),
        Inches(1.5),
        Inches(7),
        Inches(5),
        chart_data,
    )


def _add_gst_by_slab_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_GST_BY_SLAB)
    rows = data.gst_by_slab or [
        GstSlabTotal(gst_slab=0, tax_collected_paise=0, taxable_paise=0)
    ]
    chart_data = _new_category_chart_data()
    chart_data.categories = [f"{row.gst_slab}%" for row in rows]
    _add_chart_series(
        chart_data,
        "Tax collected (paise)",
        [row.tax_collected_paise for row in rows],
    )
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5),
        chart_data,
    )


def _add_stock_health_slide(presentation: Any, data: AnalysisDeckData) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_STOCK_HEALTH)
    products = data.low_stock[:8]
    if not products:
        body = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8), Inches(2))
        body.text_frame.text = "No Products at or below Reorder Level."
        return

    chart_data = _new_category_chart_data()
    chart_data.categories = [product.name for product in products]
    _add_chart_series(
        chart_data,
        "Quantity",
        [_quantity_as_float(product.quantity) for product in products],
    )
    _add_chart_series(
        chart_data,
        "Reorder Level",
        [_quantity_as_float(product.reorder_level) for product in products],
    )
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5),
        chart_data,
    )


def _add_insights_slide(presentation: Any, insights_text: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _set_title(slide, SLIDE_INSIGHTS)
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(5))
    text_frame = body.text_frame
    text_frame.word_wrap = True
    text_frame.text = insights_text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)


def _set_title(slide: object, title: str) -> None:
    shapes = getattr(slide, "shapes")
    if shapes.title is not None:
        shapes.title.text = title
    else:
        box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        box.text_frame.text = title


def _quantity_as_float(quantity: str) -> float:
    return float(Decimal(quantity))


def _new_category_chart_data() -> Any:
    return CategoryChartData()  # type: ignore[no-untyped-call]


def _add_chart_series(
    chart_data: Any,
    name: str,
    values: list[float] | list[int],
) -> None:
    chart_data.add_series(name, values)
